"""PowerPoint pitch deck generator using python-pptx."""
import os
import re
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from config import get_settings
from document_parser.section_extractor import SECTION_DISPLAY_NAMES

settings = get_settings()

# Colors
NAVY = RGBColor(0, 51, 102)
TEAL = RGBColor(0, 102, 153)
GOLD = RGBColor(204, 153, 0)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(33, 33, 33)
LIGHT_BG = RGBColor(245, 247, 250)


def generate_pptx(
    report_id: str,
    title: str,
    sections: dict,
    financial_data: dict,
    project_name: str,
) -> str:
    """Generate a consulting pitch deck and return the file path."""
    os.makedirs(settings.REPORTS_DIR, exist_ok=True)
    pptx_path = os.path.join(settings.REPORTS_DIR, f"{report_id}.pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Title Slide
    _add_title_slide(prs, project_name, title)

    # --- Table of Contents
    _add_toc_slide(prs, sections)

    # --- Section Slides (multi-slide for long content)
    for key, content in sections.items():
        display_name = SECTION_DISPLAY_NAMES.get(key, key.replace("_", " ").title())
        _add_content_slides(prs, display_name, content)

    # --- Financial Highlights
    if financial_data:
        _add_financials_slide(prs, financial_data)

    # --- Thank You Slide
    _add_closing_slide(prs, project_name)

    prs.save(pptx_path)
    return pptx_path


# ------------------------------------------------------------------ #
#  Title slide                                                        #
# ------------------------------------------------------------------ #
def _add_title_slide(prs, project_name, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Gold accent bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(3.2), Inches(13.333), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "DETAILED PROJECT REPORT"
    p.font.size = Pt(40)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = project_name.upper()
    p2.font.size = Pt(28)
    p2.font.color.rgb = GOLD
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER

    # Date
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = f"Prepared: {datetime.now().strftime('%B %Y')} | Confidential"
    p3.font.size = Pt(14)
    p3.font.color.rgb = WHITE
    p3.alignment = PP_ALIGN.CENTER


# ------------------------------------------------------------------ #
#  Table of contents                                                  #
# ------------------------------------------------------------------ #
def _add_toc_slide(prs, sections):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_header(slide, "TABLE OF CONTENTS")

    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, key in enumerate(list(sections.keys())[:13], 1):
        display = SECTION_DISPLAY_NAMES.get(key, key.replace("_", " ").title())
        p = tf.add_paragraph() if i > 1 else tf.paragraphs[0]
        p.text = f"{i}.  {display}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.space_after = Pt(4)

    remaining = list(sections.keys())[13:]
    if remaining:
        txBox2 = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, key in enumerate(remaining, 14):
            display = SECTION_DISPLAY_NAMES.get(key, key.replace("_", " ").title())
            p = tf2.add_paragraph() if i > 14 else tf2.paragraphs[0]
            p.text = f"{i}.  {display}"
            p.font.size = Pt(14)
            p.font.color.rgb = DARK
            p.space_after = Pt(4)


# ------------------------------------------------------------------ #
#  Content slides (multi-slide, proper tables)                        #
# ------------------------------------------------------------------ #
def _add_content_slides(prs, title, content):
    """Add one or more content slides for a section.
    Long content is split across slides. Tables render as native PPTX tables."""
    content = content or ""
    lines = content.split("\n")

    # Parse lines into structured chunks: text blocks and table blocks
    chunks = []
    text_buf = []
    table_buf = []
    in_table = False

    for line in lines:
        stripped = line.strip()
        is_tbl = stripped.startswith("|") and "|" in stripped[1:]
        is_sep = bool(re.match(r"^\|[\s\-:|]+\|", stripped))

        if is_tbl or is_sep:
            if not in_table:
                if text_buf:
                    chunks.append(("text", text_buf))
                    text_buf = []
                in_table = True
            if not is_sep:
                table_buf.append(stripped)
        else:
            if in_table:
                if table_buf:
                    chunks.append(("table", table_buf))
                    table_buf = []
                in_table = False
            if stripped:
                text_buf.append(stripped)

    if table_buf:
        chunks.append(("table", table_buf))
    if text_buf:
        chunks.append(("text", text_buf))

    # Paginate into slide-pages
    MAX_LINES = 16
    pages = []          # list of list-of-(type, data)
    cur_page = []
    cur_count = 0

    for ctype, cdata in chunks:
        if ctype == "table":
            need = len(cdata) + 2
            if cur_count + need > MAX_LINES and cur_page:
                pages.append(cur_page)
                cur_page = []
                cur_count = 0
            cur_page.append(("table", cdata))
            cur_count += need
        else:
            for ln in cdata:
                if cur_count >= MAX_LINES:
                    pages.append(cur_page)
                    cur_page = []
                    cur_count = 0
                cur_page.append(("line", ln))
                cur_count += 1

    if cur_page:
        pages.append(cur_page)

    if not pages:
        pages = [[("line", "Content not available.")]]

    # Render each page
    for pidx, page in enumerate(pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        stitle = title if pidx == 0 else f"{title} (contd.)"
        _add_slide_header(slide, stitle)

        y = Inches(1.5)
        for itype, idata in page:
            if itype == "table":
                y = _render_table(slide, idata, y)
            else:
                y = _render_line(slide, idata, y)


def _render_line(slide, line, y):
    """Render one line of markdown text as a PPTX text box."""
    clean = line.strip()
    if not clean:
        return y

    is_heading = clean.startswith("#")
    is_bullet = clean.startswith("- ") or clean.startswith("* ") or clean.startswith("  -")
    num_match = re.match(r"^(\d+)\.\s+(.+)", clean)
    is_emoji = any(clean.startswith(e) for e in ["  ", "  ", "  ", "  "])

    if is_heading:
        level = len(re.match(r"^#+", clean).group())
        clean = clean.lstrip("#").strip()
        font_size = {1: 18, 2: 16, 3: 14, 4: 13}.get(level, 12)
        bold = True
        color = NAVY
        indent = 0.8
    elif is_bullet:
        clean = clean.lstrip("-*  ").strip()
        clean = f"  {clean}"
        font_size = 11
        bold = False
        color = DARK
        indent = 1.0
    elif num_match:
        clean = f"{num_match.group(1)}.  {num_match.group(2)}"
        font_size = 11
        bold = False
        color = DARK
        indent = 1.0
    else:
        font_size = 11
        bold = False
        color = DARK
        indent = 0.8

    # Strip markdown bold/italic
    clean = clean.replace("**", "").replace("*", "")
    # Replace arrow markers
    clean = clean.replace(" -> ", " > ").replace("->", " > ")

    txBox = slide.shapes.add_textbox(Inches(indent), y, Inches(11.5), Inches(0.35))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = clean
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_after = Pt(2)

    return y + Inches(0.28)


def _render_table(slide, table_lines, y):
    """Render a markdown-style table as a native PPTX table."""
    if not table_lines:
        return y

    rows_data = []
    for line in table_lines:
        cells = [c.strip() for c in line.split("|")]
        cells = [c for c in cells if c != ""]
        if cells:
            rows_data.append(cells)

    if not rows_data:
        return y

    num_cols = max(len(r) for r in rows_data)
    num_rows = len(rows_data)

    tw = Inches(11.5)
    rh = Inches(0.32)

    ts = slide.shapes.add_table(num_rows, num_cols, Inches(0.8), y, tw, rh * num_rows)
    table = ts.table

    for ri, row in enumerate(rows_data):
        for ci in range(num_cols):
            cell = table.cell(ri, ci)
            txt = row[ci].replace("**", "").replace("*", "") if ci < len(row) else ""
            cell.text = txt
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)

            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
                p.font.color.rgb = WHITE
                p.font.bold = True
            else:
                p.font.color.rgb = DARK
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_BG

    return y + rh * num_rows + Inches(0.15)


# ------------------------------------------------------------------ #
#  Financial highlights                                               #
# ------------------------------------------------------------------ #
def _add_financials_slide(prs, financial_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_header(slide, "FINANCIAL HIGHLIGHTS")

    pc = financial_data.get("project_cost", {})
    mof = financial_data.get("means_of_finance", {})
    bep = financial_data.get("breakeven", {})
    projections = financial_data.get("revenue_projections", [])

    highlights = [
        ("Total Project Cost", _format_inr(pc.get("total", 0))),
        ("Promoter's Contribution", _format_inr(mof.get("promoter_contribution", 0))),
        ("Term Loan", _format_inr(mof.get("term_loan", 0))),
        ("Debt-Equity Ratio", f"{mof.get('debt_equity_ratio', 0)}:1"),
        ("Break-Even Point", f"{bep.get('bep_percentage', 0)}% of capacity"),
    ]

    if projections:
        yr1 = projections[0]
        yr5 = projections[-1]
        highlights.append(("Year 1 Revenue", _format_inr(yr1.get("revenue", 0))))
        highlights.append(("Year 1 Net Profit", _format_inr(yr1.get("net_profit", 0))))
        highlights.append(("Year 5 Revenue", _format_inr(yr5.get("revenue", 0))))

    for i, (label, value) in enumerate(highlights[:8]):
        col = i % 4
        row = i // 4
        x = Inches(0.8 + col * 3.1)
        y_pos = Inches(1.8 + row * 2.5)

        shape = slide.shapes.add_shape(1, x, y_pos, Inches(2.8), Inches(2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BG
        shape.line.color.rgb = TEAL
        shape.line.width = Pt(1)

        txBox = slide.shapes.add_textbox(x + Inches(0.15), y_pos + Inches(0.3), Inches(2.5), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(20)
        p.font.color.rgb = NAVY
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        txBox2 = slide.shapes.add_textbox(x + Inches(0.15), y_pos + Inches(1.2), Inches(2.5), Inches(0.6))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = label
        p2.font.size = Pt(11)
        p2.font.color.rgb = DARK
        p2.alignment = PP_ALIGN.CENTER


# ------------------------------------------------------------------ #
#  Closing slide                                                      #
# ------------------------------------------------------------------ #
def _add_closing_slide(prs, project_name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "THANK YOU"
    p.font.size = Pt(44)
    p.font.color.rgb = GOLD
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = project_name
    p2.font.size = Pt(18)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = f"Prepared by {project_name}"
    p3.font.size = Pt(12)
    p3.font.color.rgb = WHITE
    p3.alignment = PP_ALIGN.CENTER


# ------------------------------------------------------------------ #
#  Helpers                                                            #
# ------------------------------------------------------------------ #
def _add_slide_header(slide, title):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.bold = True

    line = slide.shapes.add_shape(1, Inches(0), Inches(1.2), Inches(13.333), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()


def _format_inr(amount):
    try:
        amount = float(amount)
    except (ValueError, TypeError):
        return str(amount)
    if amount >= 10000000:
        return f"Rs. {amount/10000000:.2f} Cr"
    elif amount >= 100000:
        return f"Rs. {amount/100000:.2f} L"
    else:
        return f"Rs. {amount:,.0f}"
